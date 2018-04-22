from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import *
import sys

f = open(sys.argv[1],"r")
lines = str(f.read()).split('\n')

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

for line in lines:
    slide = prs.slides.add_slide(blank_slide_layout)

    thr = (2,6,18)

    f_size = 60 if len(line)==0 else 600/len(line) if thr[0]<=len(line)<=thr[1] \
        else 2000/len(line) if len(line)>=thr[2] else 600/thr[0] if len(line)<thr[0] else 600/thr[1]

    txBox = slide.shapes.add_textbox(0, Cm(-f_size/45), Cm(25.4), Cm(19.05))
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = line
    p.font.bold = True
    p.font.size = Pt(f_size)
    p.alignment = PP_ALIGN.CENTER

    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.fit_text(font_family='Calibri', max_size=f_size, bold=True, italic=False, font_file=None)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

name = str(sys.argv[1].split('.')[0])
prs.save(name+'.pptx')